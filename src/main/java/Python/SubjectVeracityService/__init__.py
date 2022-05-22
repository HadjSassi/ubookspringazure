import nltk
from nltk.corpus import stopwords
import numpy as np
from docx import Document
import PyPDF2
from pptx import Presentation
from joblib import dump, load

matiere = ''
file = ''

def docx_to_dict(file_name):
    docx_dict = {}
    document = Document(file_name)
    indx = 0
    for para in document.paragraphs:
        indx += 1
        if (len(para.text) > 0):
            docx_dict[indx] = para.text
    return docx_dict


def texting(doc):
    # print(doc)
    df = docx_to_dict(doc)
    txt = ""
    for pp in df.values():
        txt += pp + "\n"
    return txt


def pdfinf(chaine):
    # print(chaine)
    tt = ""
    pdf = PyPDF2.PdfFileReader(chaine)
    for i in range(pdf.numPages):
        page_1_object = pdf.getPage(i)
        page_1_text = page_1_object.extractText()
        tt += page_1_text
    return tt


def pptxing(chaine):
    # print(chaine)
    tt = ""
    prs = Presentation(chaine)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                tt += shape.text
    return tt


def clearExtraSpaces(text):
    text = text.strip()
    text = text.replace("\n", "")
    text = text.replace("  ", "")
    text = text.replace("\n", "")
    text = text.replace("\t", " ")
    return text


def tokenisation(text):
    final_stopwords_list = stopwords.words('english') + stopwords.words('french')
    final_stopwords_list.extend(["'", ":", "(", ")", "?", "!", "[", "]", ",", "’", ".", "»", ">", "<", "«", "%"])
    tokens = nltk.word_tokenize(text)
    token = []
    for tok in tokens:
        tok = tok.lower()
        tok = tok.replace("à", "a")
        tok = tok.replace("é", "e")
        tok = tok.replace("é", "e")
        tok = tok.replace("î", "i")
        tok = tok.replace("ô", "o")
        tok = tok.replace("û", "u")
        tok = tok.replace("â", "a")
        if tok not in final_stopwords_list:
            token.append(tok)
    return token


def extractionText(file):
    txt = ''
    if file.endswith(".pdf"):
        txt = pdfinf(file)
    elif file.endswith(".pptx"):
        txt = pptxing(file)
    elif file.endswith(".docx"):
        txt = texting(file)
    else:
        print("File unknown type")
    return txt


def treat(file):
    txt = extractionText(file)
    tokens = tokenisation(txt)
    return tokens



def find(path, file=file):
    tokens = treat(file)
    with open(f"{path}.txt", "r", encoding="utf8") as f:
        x = f.read()
    bagOfWords = x.split("\n")
    bagOfWords.pop()
    occurence = [0 for j in bagOfWords]
    for tt in tokens:
        if (tt in bagOfWords):
            k = bagOfWords.index(tt)
            occurence[k] += 1
    dataSets = np.array([
        occurence
    ])

    model = load(f"{path}.joblib")
    y = model.predict(dataSets)
    return y[0]


def guess(file):
    isSc = find("src/main/java/Python/SubjectVeracityService/IsScientific",file)
    if (isSc == "science"):
        isScExp = find("src/main/java/Python/SubjectVeracityService/IsScExp",file)
        if (isScExp == "mpci"):
            mpci = find("src/main/java/Python/SubjectVeracityService/IsChemistry",file)
            if (mpci == "mpi"):
                mpi = find("src/main/java/Python/SubjectVeracityService/IsInfo",file)
                if(mpi == "mp"):
                    matiere = find("src/main/java/Python/SubjectVeracityService/IsPhysics",file)
                else:
                    matiere = mpi
            else:
                matiere = mpci
        else:
            matiere = find("src/main/java/Python/SubjectVeracityService/IsSante",file)
    else:
        matiere = find("src/main/java/Python/SubjectVeracityService/IsLiterature",file)

    if (matiere == "type 1"):
        matiere = "Sante"
    elif (matiere == "Here is the type 2"):
        matiere = "SVT"
    elif (matiere == "math"):
        matiere = "Maths"
    elif (matiere == "phy"):
        matiere = "Physics"
    elif (matiere == "./lanTxt"):
        matiere = "Literature"
    elif (matiere == "./manTxt"):
        matiere = "Management"
    elif (matiere == "chimie"):
        matiere = "Chemistry"

    return matiere

class SubjectVeracityService :
    matiere = ""
    def __init__(self, file):
        self.matiere = guess(file)
    def getSubject(self):
        return self.matiere
